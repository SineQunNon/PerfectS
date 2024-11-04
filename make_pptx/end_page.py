from pptx.util import Pt
from pptx.enum.text import PP_ALIGN,MSO_ANCHOR

def make_end_page(prs):
    prs.slides.add_slide(prs.slide_layouts[5])
    prs.slides.add_slide(prs.slide_layouts[6])