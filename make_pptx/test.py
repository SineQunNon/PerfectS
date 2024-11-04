from pptx import Presentation
from make_pptx.get_template import get_template
from make_pptx.title_page import make_title_page


prs = Presentation("/Users/sinequanon/Documents/PerfectS_beta/make_pptx/template/template_final.pptx")

make_title_page(prs)

slide = prs.slides.add_slide(prs.slide_layouts[2])

title_shape = slide.shapes.title
if title_shape:
    title_shape.text = "Signalment"
else:
    pass
titlenum_shape = slide.shapes[2]
titlenum_shape.text = "01"
body_shape = slide.shapes[1]
body_shape.text = """• Name : 먼지
• Breed : 코리안 숏헤어
• Species : 고양이
• Age : 2y 6m
• Sex : CM
• BW : 6.2 Kg
"""


prs.save("/Users/sinequanon/Documents/PerfectS_beta/data/pptx_output/test.pptx")