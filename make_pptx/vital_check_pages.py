from Rag.loader.pdf_loader import load_pdf
from Rag.loader.main import get_vital_check
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # GUI 백엔드가 아닌 Agg 백엔드로 변경
from pptx.util import Inches, Cm, Pt
import os
import pandas as pd
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def plot_vital_check_graph(vital_check_data_list):
    # vital_check_data_list에서 날짜와 BW(Kg) 추출
    dates = vital_check_data_list["날짜"].tolist()
    bw_values = vital_check_data_list["BW(Kg)"].tolist()

    # 추출된 데이터 출력
    #print("Dates:", dates)
    #print("BW Values:", bw_values)
    min_bw = float(min(bw_values)) - 1.0
    max_bw = float(max(bw_values)) + 1.0
    
    if not dates or not bw_values:
        #print("날짜와 BW(Kg) 값을 찾을 수 없습니다.")
        return

    #bw 데이터 표시
    bw_values = [float(value) for value in bw_values]

    # x값 정렬
    indices = list(range(1, len(dates) + 1))

    # 그래프
    plt.figure(figsize=(12, 6))
    plt.plot(dates, bw_values, marker='o', label='BW (Kg)')
    plt.xlabel('date')
    plt.ylabel('BW(Kg)')
    plt.ylim(min_bw, max_bw)
    plt.title('BW(Kg) Over Time')
    plt.legend()
    plt.grid(True)

    # 이미지 저장
    graph_path = 'vital_check_graph.png'
    plt.tight_layout()
    plt.savefig(graph_path)
    plt.close()

    return graph_path

#vital check 표 추가
def add_vital_check_table_to_slide(prs, num, headers, vital_check_table):
    slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(slide_layout)

    if num < 10:
        slide.shapes[1].text = "0" + str(num)
    else:
        slide.shapes[1].text = str(num)

    title = slide.shapes.title
    if title:
        title.text = "Vital Check Table"
    
    # 테이블을 삽입할 위치와 크기 설정
    x, y, cx, cy = Inches(0.88), Inches(1.6), Inches(11.68), Inches(2)

    rows, cols = vital_check_table.shape[0] + 1, vital_check_table.shape[1]
    table = slide.shapes.add_table(rows, cols, x, y, cx, cy).table

    table.columns[0].width = Inches(2)  # 첫 번째 열 너비
    for col_idx in range(1, cols):
        table.columns[col_idx].width = Inches(1.6)

    # 헤더 스타일 설정
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(227, 131, 85)  # 배경색 설정 (주황색)
        cell.text_frame.paragraphs[0].font.name = "Noto Sans KR"
        cell.text_frame.paragraphs[0].font.bold = True     # 굵게
        cell.text_frame.paragraphs[0].font.size = Pt(18)   # 글꼴 크기 설정
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 글자 색상 흰색
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # 가운데 정렬
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # 중간 정렬

    # 데이터 셀 스타일 설정
    for row_idx, row in vital_check_table.iterrows():
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.fill.solid()
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(230, 230, 230)  # 배경색 회색 (짝수 행)
            else:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 배경색 흰색 (홀수 행)
            
            cell.text_frame.paragraphs[0].font.size = Pt(18)  # 글꼴 크기
            cell.text_frame.paragraphs[0].font.name = "Noto Sans KR"
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # 가운데 정렬
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # 중간 정렬

def make_vital_check_page(prs, pdf_path):
    """ vital check 그래프 삽입 """
    pages = load_pdf(pdf_path)
    vital_check_table = get_vital_check(pages)

    graph_path = plot_vital_check_graph(vital_check_table)

    idx = len(prs.slides)
    slide = prs.slides[idx-1]
    num = int(slide.shapes[2].text) + 1

    #print(num)

    if graph_path:
        slide_layout = prs.slide_layouts[4]
        graph_slide = prs.slides.add_slide(slide_layout)

        titlenum_shape = graph_slide.shapes[1]
        if num < 10:
            titlenum_shape.text = "0" + str(num)
        else:
            titlenum_shape.text = str(num)

        title = graph_slide.shapes.title
        if title:
            title.text = "Vital Check Graph"
        
        graph_slide.shapes.add_picture(graph_path, Inches(1), Inches(1.7), width=Inches(6.5), height=Inches(3.65))

        # 이미지 파일 제거
        try:
            if os.path.exists(graph_path):
                os.remove(graph_path)
                #print(f"Temporary image file {graph_path} removed.")
        except Exception as e:
            print(f"Error removing file {graph_path}: {e}")

    
    """ vital check 테이블 삽입 """

    vital_headers = ['날짜', '시간', 'BW(Kg)', 'BT(C)', 'BP(mmHg)', 'HR(/min)', 'Sign']

    #print(vital_check_table)
    
    add_vital_check_table_to_slide(prs, num+1,vital_headers, vital_check_table)