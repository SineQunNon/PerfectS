from Rag.loader.main import get_pptx_hct_tables, get_pptx_AnionGap_and_other_tables
from Rag.llm.get_ppt_input import get_lab_table_summary_input
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN,MSO_ANCHOR
from pptx.util import Pt
import logging

def make_lab_hct_page(prs, pdf_path):
    df = get_pptx_hct_tables(pdf_path)
    #print(df)
    if not df.empty:
        idx = len(prs.slides)
        t_slide = prs.slides[idx-1]
        num = int(t_slide.shapes[1].text) + 1
        #print(num)

        slide_layout = prs.slide_layouts[3]  # 빈 슬라이드
        slide = prs.slides.add_slide(slide_layout)

        titlenum_shape = slide.shapes[1]
        if num < 10:
            titlenum_shape.text = "0" + str(num)
        else:
            titlenum_shape.text = str(num)

        title = slide.shapes.title
        if title:
            title.text = "Lab Table(혈구 검사)"
        
        answer = get_lab_table_summary_input(str(df))
        body_shape = slide.shapes[2]
        body_shape.text = f"{answer}(ChatGPT요약)"

        # 테이블을 삽입할 위치와 크기 설정
        x, y, cx, cy = Inches(0.88), Inches(1.6), Inches(11.68), Inches(2)

        # 날짜 컬럼을 헤더로 사용하므로, 중복된 날짜는 제거하여 헤더로 생성
        dates = df['date'].unique()

        # 행은 '검사명'이기 때문에 3개(검사명 + 범위), 열은 날짜 개수 + 1 (첫 열은 검사명, 범위 등)
        rows = len(df['검사명'].unique()) + 1  # 검사명 행과 범위 행
        cols = len(dates) + 1

        table = slide.shapes.add_table(rows, cols, x, y, cx, cy).table

        # 첫 번째 열 너비를 넓게 설정
        table.columns[0].width = Inches(1.6)
        # 0번째 열의 셀 색상 변경
        cell = table.cell(0, 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(227, 131, 85)  # 주황색
        # 헤더에 날짜를 추가
        for col_idx, date in enumerate(dates, start=1):
            cell = table.cell(0, col_idx)
            cell.text = str(date)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(227, 131, 85)  # 주황색
            cell.text_frame.paragraphs[0].font.name = "Noto Sans KR"
            cell.text_frame.paragraphs[0].font.bold = True           
            cell.text_frame.paragraphs[0].font.size = Pt(18)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 흰색
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # 중간 정렬


        # 검사명, 범위 추가 (첫 번째 열에 검사명과 범위 표시)
        test_names = df['검사명'].unique()
        for row_idx, test_name in enumerate(test_names):
            # 검사명
            cell = table.cell(row_idx + 1, 0)
            cell.text = test_name
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(227, 131, 85)  
            cell.text_frame.paragraphs[0].font.name = "Noto Sans KR"
            cell.text_frame.paragraphs[0].font.size = Pt(18)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 흰색
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # 중간 정렬

        # 결과값 및 범위 추가
        for col_idx, date in enumerate(dates, start=1):
            filtered_data = df[df['date'] == date]
            for row_idx, test_name in enumerate(test_names):
                result_row = filtered_data[filtered_data['검사명'] == test_name]

                # 검사 결과 값 추가
                result_cell = table.cell(row_idx + 1, col_idx)
                value = result_row['결과값'].values[0]
                status = result_row['status'].values[0]
                if status == '▲':  # 상태 표시
                    result_cell.text = f"{value}({status})"
                    result_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)  # 빨간색
                elif status == '▼':
                    result_cell.text = f"{value}({status})"
                    result_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 255)  # 파란색

                else:
                    result_cell.text = str(value)

                # 셀 스타일 적용
                result_cell.fill.solid()
                result_cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 흰색
                result_cell.text_frame.paragraphs[0].font.name = "Noto Sans KR"
                result_cell.text_frame.paragraphs[0].font.size = Pt(18)
                result_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                result_cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # 중간 정렬

    else:
        pass
        #logging.info("lab table is empty")

def get_table_str(table_list):
    table_str = ""
    for data in table_list:
        table_str += str(data['table'])
        table_str += '\n'
    return table_str

def make_lab_aniongap_other_page(prs, pdf_path):
    anion_table_list, other_table_list = get_pptx_AnionGap_and_other_tables(pdf_path)
    
    # 가스 분석 표 생성
    if len(anion_table_list) != 0:
        
        idx = len(prs.slides)
        t_slide = prs.slides[idx-1]
        num = int(t_slide.shapes[1].text) + 1# 슬라이드 타이틀 번호
            
        slide_layout = prs.slide_layouts[3]
        slide = prs.slides.add_slide(slide_layout)

        titlenum_shape = slide.shapes[1]
        if num < 10:
            titlenum_shape.text = "0" + str(num)
        else:
            titlenum_shape.text = str(num)

        title = slide.shapes.title
        if title:
            title.text = "Lab Table(가스 분석)"
        
        table_str = get_table_str(anion_table_list)
        answer = get_lab_table_summary_input(table_str)
        body_shape = slide.shapes[2]
        body_shape.text = f"{answer}(ChatGPT요약)"

        x1, y1, cx1, cy1 = Inches(0.80), Inches(1.6), Inches(5.5), Inches(2)
        x2, y2, cx2, cy2 = Inches(7.0), Inches(1.6), Inches(5.5), Inches(2)
        
        # 검사명 추출
        test_names = anion_table_list[0]['table']['검사명']
        test_length = len(test_names)
        
        # 테이블 행 길이 판단
        if test_length % 2 == 0:
            first_table_length = test_length // 2 + 1
            second_table_length = test_length // 2 + 1
        else:
            first_table_length = test_length // 2  + 2
            second_table_length = test_length // 2 + 1
        
        cols = len(anion_table_list) + 1
        
        # 테이블 생성
        first_table = slide.shapes.add_table(first_table_length, cols, x1, y1, cx1, cy1).table
        second_table = slide.shapes.add_table(second_table_length, cols, x2, y2, cx2, cy2).table
        
        first_column = first_table.columns[0]
        first_column.width = Inches(1.0)
        first_column = second_table.columns[0]
        first_column.width = Inches(1.0)
        # 각 테이블 폰트 사이즈 변경
        for row in first_table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(9)  # 폰트 크기를 8로 설정
            #row.height = Inches(0.2)  # 각 행의 높이를 원하는 값으로 설정
        
        for row in second_table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(9)  # 폰트 크기를 8로 설정
            #row.height = Inches(0.2)  # 각 행의 높이를 원하는 값으로 설정
        
        
        # 첫 번째 테이블 검사명 넣기
        for idx, row in enumerate(first_table.rows, start=1):
            if idx==first_table_length:
                break
            
            cell = first_table.cell(idx, 0)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(227,131,85)
            cell.text = test_names[idx-1]
            cell.text_frame.paragraphs[0].font.size = Pt(8)  # 텍스트 크기 설정
            #cell.text_frame.paragraphs[0].font.bold = True   # 굵게 설정 (필요 시)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # 가운데 정렬 (필요 시)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
        # 두 번째 테이블 검사명 넣기
        for idx, row in enumerate(second_table.rows, start=1):
            if idx == second_table_length:
                break
            cell = second_table.cell(idx, 0)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(227,131,85)
            cell.text = test_names[idx+first_table_length-2]
            cell.text_frame.paragraphs[0].font.size = Pt(8)  # 텍스트 크기 설정
            #cell.text_frame.paragraphs[0].font.bold = True   # 굵게 설정 (필요 시)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # 가운데 정렬 (필요 시)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
        
        first_cell = first_table.cell(0, 0)
        first_cell.fill.solid()
        first_cell.fill.fore_color.rgb = RGBColor(227,131, 85)
        
        second_cell = second_table.cell(0, 0)
        second_cell.fill.solid()
        second_cell.fill.fore_color.rgb = RGBColor(227,131, 85)
        
        dates = []
        for df in anion_table_list:
            date = df['date']
            dates.append(date)
        
        # 0 번째 행 날짜 채우기
        for col_idx, date in enumerate(dates, start=1):
            # 첫 번째 테이블
            first_cell = first_table.cell(0, col_idx)
            first_cell.text = str(date)
            first_cell.fill.solid()
            first_cell.fill.fore_color.rgb = RGBColor(227,131,85)
            first_cell.text_frame.paragraphs[0].font.name = "Noto Sans KR"
            first_cell.text_frame.paragraphs[0].font.bold = True
            first_cell.text_frame.paragraphs[0].font.size = Pt(10)
            first_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
            first_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            first_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # 두 번째 테이블
            second_cell = second_table.cell(0, col_idx)
            second_cell.text = str(date)
            second_cell.fill.solid()
            second_cell.fill.fore_color.rgb = RGBColor(227,131,85)
            second_cell.text_frame.paragraphs[0].font.name = "Noto Sans KR"
            second_cell.text_frame.paragraphs[0].font.bold = True
            second_cell.text_frame.paragraphs[0].font.size = Pt(8)
            second_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
            second_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            second_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            #print(anion_table_list)
            for col in range(1, len(dates)+1):
                selected_table = anion_table_list[col-1]['table']
                for row in range(1, first_table_length):
                    cell = first_table.cell(row, col)
                    text = str(selected_table.iloc[row-1]['결과값'])
                    status = selected_table.iloc[row-1]['status']
                    if status == '▲':
                        cell.text = f"{text}({status})"
                        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
                    elif status == '▼':
                        cell.text = f"{text}({status})"
                        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 255)
                    else:
                        cell.text = text
                    cell.text_frame.paragraphs[0].font.name = "Noto Sans KR"
                    cell.text_frame.paragraphs[0].font.size = Pt(8)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 흰색
            
            for col in range(1, len(dates)+1):
                selected_table = anion_table_list[col-1]['table']
                for row in range(1, second_table_length):
                    idx = row + first_table_length - 1
                    cell = second_table.cell(row, col)
                    text = str(selected_table.iloc[idx-1]['결과값'])
                    status = selected_table.iloc[idx-1]['status']
                    if status == '▲':
                        cell.text = f"{text}({status})"
                        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
                    elif status == '▼':
                        cell.text = f"{text}({status})"
                        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 255)
                    else:
                        cell.text = text
                    cell.text_frame.paragraphs[0].font.name = "Noto Sans KR"
                    cell.text_frame.paragraphs[0].font.size = Pt(8)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 흰색
    
    # 나머지 테이블
    # if len(other_table_list) != 0:
    #     idx = len(prs.slides)
    #     t_slide = prs.slides[idx-1]
    #     num = int(t_slide.shapes[1].text) + 1 # 슬라이드 타이틀 번호
            
    #     slide_layout = prs.slide_layouts[3]

        # for df in other_df:
        #     slide = prs.slides.add_slide(slide_layout)

        #     titlenum_shape = slide.shapes[1]
        #     if num < 10:
        #         titlenum_shape.text = "0" + str(num)
        #     else:
        #         titlenum_shape.text = str(num)
            
        #     title = slide.shapes.title
        #     if title:
        #         title.text = "Lab Table(others)"