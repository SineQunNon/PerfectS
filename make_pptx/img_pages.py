import tempfile
from PIL import Image
import fitz
import cv2
import io
import numpy as np
import os
from pptx.util import Inches, Cm
import logging

def remove_white_background(image):
    # 이미지를 numpy 배열로 변환
    open_cv_image = np.array(image)
    if open_cv_image.shape[2] == 4:
        open_cv_image = cv2.cvtColor(open_cv_image, cv2.COLOR_RGBA2RGB)
    # 이미지를 그레이스케일로 변환
    gray = cv2.cvtColor(open_cv_image, cv2.COLOR_RGB2GRAY)
    # 흰색 배경을 바이너리 이미지로 변환 (임계값 240 이상을 흰색으로 간주)
    _, binary = cv2.threshold(gray, 240, 255, cv2.THRESH_BINARY_INV)
    # 바이너리 이미지에서 외곽선을 찾음
    contours, _ = cv2.findContours(binary, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    # 외곽선이 있는 경우 가장 큰 외곽선을 기준으로 이미지 자름
    if contours:
        x, y, w, h = cv2.boundingRect(contours[0])
        for contour in contours:
            x_, y_, w_, h_ = cv2.boundingRect(contour)
            if w_ * h_ > w * h:
                x, y, w, h = x_, y_, w_, h_
        # 외곽선을 포함하는 사각형 영역을 자름
        cropped_image = open_cv_image[y:y+h, x:x+w]
        return Image.fromarray(cropped_image)
    return image

def extract_images_from_pdf(pdf_path):
    # PDF 문서 열기
    pdf_document = fitz.open(pdf_path)
    images = []
    for page_num in range(len(pdf_document)):
        page = pdf_document[page_num]
        # 페이지에서 이미지 목록 가져오기
        image_list = page.get_images(full=True)
        # 각 이미지 추출
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            # 이미지 데이터를 PIL 이미지로 변환
            image = Image.open(io.BytesIO(image_bytes))
            # 흰색 배경 제거
            image = remove_white_background(image)
            # 이미지를 리스트에 추가
            images.append(image)
            #print(f"Extracted image {img_index} from page {page_num}, size: {image.size}, mode: {image.mode}")

    # Debugging: Print the number of extracted images
    logging.info(f"Total number of extracted images: {len(images)}")

    return images

# 차트 이미지 페이지 1개 생성
def add_new_slide_with_chart_image(prs, num, img_left, img_right=None):

    # 슬라이드 3번 템플릿 가져오기
    slide_layout = prs.slide_layouts[2]

    # 페이지 생성
    new_slide = prs.slides.add_slide(slide_layout)
    
    titlenum_shape = new_slide.shapes[2]
    if num < 10:
        titlenum_shape.text = "0" + str(num)
    else:
        titlenum_shape.text = str(num)
    # 타이틀 추가
    title = new_slide.shapes.title

    if title:
        title.text = "Chart Image"

    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img_left:
        # 첫 번째 이미지 임시저장
        img_left.save(tmp_img_left, format="PNG")
        tmp_img_left_path = tmp_img_left.name
        #print(f"Temporary Saving to {tmp_img_left_path}")
    
    if img_left:
        try:
            width, height = img_left.size
            max_width = Inches(5)
            max_height = Inches(4)

            # 크기 비율 계산
            ratio = min(max_width / width, max_height/height)
            new_width = width * ratio
            new_height = height * ratio

            new_slide.shapes.add_picture(tmp_img_left_path, Inches(1), Inches(1.6), width=new_width, height=new_height)
            #print("successfully added left image to slide")
        except Exception as e:
            logging.error(f"Error adding left image to slide: {e}")

    if img_right:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img_right:
            img_right.save(tmp_img_right, format="PNG")
            tmp_img_right_path = tmp_img_right.name
            #print(f"Temporary Saving to {tmp_img_right_path}")
        
        if os.path.exists(tmp_img_right_path):
            try:
                width, height = img_right.size
                max_width = Inches(5)
                max_height = Inches(4)

                # 크기 비율 계산
                ratio = min(max_width / width, max_height/height)
                new_width = width * ratio
                new_height = height * ratio

                new_slide.shapes.add_picture(tmp_img_right_path, Inches(7), Inches(1.6), width=new_width, height=new_height)
                #print("successfully added right image to slide")
            except Exception as e:
                logging.error(f"Error adding right image to slide: {e}")


def make_chart_img_pages(prs, pdf_path):
    images = extract_images_from_pdf(pdf_path)

    idx = len(prs.slides)
    slide = prs.slides[idx-1]
    t_num = int(slide.shapes[0].text) + 1

    for i in range(0, len(images), 2):
        add_new_slide_with_chart_image(prs, t_num, images[i], images[i+1] if i+1 < len(images) else None)

    if len(images) % 2 != 0:
        add_new_slide_with_chart_image(prs, t_num, images[-1])