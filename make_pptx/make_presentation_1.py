import json
import fitz  # PyMuPDF
import io
from PIL import Image
import os
import cv2
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.dml.color import RGBColor
from pptx.shapes.group import GroupShape
from pptx.shapes.connector import Connector
import sys
#sys.path.append('/home/user/PerfectS/extract_data/')
from extract_data.extract_pdf_data import main
import tempfile
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from loader.pdf_loader import load_pdf
from loader.main import *
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR


from dotenv import load_dotenv
load_dotenv()
# 이미지에서 흰색 배경 제거
def remove_white_background(image):
    # 이미지를 numpy 배열로 변환
    open_cv_image = np.array(image)
    if open_cv_image.shape[2] == 4:
        open_cv_image = cv2.cvtColor(open_cv_image, cv2.COLOR_RGBA2RGB)
    # 이미지를 그레이스케일로 변환
    gray = cv2.cvtColor(open_cv_image, cv2.COLOR_RGB2GRAY)
    # 흰색 배경을 바이너리 이미지로 변환 (임계값 240 이상을 흰색으로 간주)
    _, binary = cv2.threshold(gray, 240, 255, cv2.THRESH_BINARY_INV)
    # 바이너리 이미지에서 외곽선을 찾음
    contours, _ = cv2.findContours(binary, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    # 외곽선이 있는 경우 가장 큰 외곽선을 기준으로 이미지 자름
    if contours:
        x, y, w, h = cv2.boundingRect(contours[0])
        for contour in contours:
            x_, y_, w_, h_ = cv2.boundingRect(contour)
            if w_ * h_ > w * h:
                x, y, w, h = x_, y_, w_, h_
        # 외곽선을 포함하는 사각형 영역을 자름
        cropped_image = open_cv_image[y:y+h, x:x+w]
        return Image.fromarray(cropped_image)
    return image

def extract_images_from_pdf(pdf_path):
    # PDF 문서 열기
    pdf_document = fitz.open(pdf_path)
    images = []
    for page_num in range(len(pdf_document)):
        page = pdf_document[page_num]
        # 페이지에서 이미지 목록 가져오기
        image_list = page.get_images(full=True)
        # 각 이미지 추출
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            # 이미지 데이터를 PIL 이미지로 변환
            image = Image.open(io.BytesIO(image_bytes))
            # 흰색 배경 제거
            image = remove_white_background(image)
            # 이미지를 리스트에 추가
            images.append(image)
            print(f"Extracted image {img_index} from page {page_num}, size: {image.size}, mode: {image.mode}")

    # Debugging: Print the number of extracted images
    print(f"Total number of extracted images: {len(images)}")

    # 첫 5개의 이미지를 파일로 저장하고 출력
    for i, img in enumerate(images[:5]):  #
        img_path = f'extracted_image_{i+1}.png'
        img.save(img_path)
        print(f"Saved {img_path}")

    return images

def select_shape_by_text(slide, text):
    # 슬라이드의 모든 도형 순회
    for shape in slide.shapes:
        # 도형이 텍스트 프레임을 가지고 있는지 확인
        if shape.has_text_frame and text in shape.text:
            return shape
    return None

def select_shapes_by_texts(slide, texts):
    # 일치하는 도형을 저장할 리스트 초기화
    shapes_with_texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            # 도형의 텍스트가 주어진 텍스트 목록 중 하나를 포함하는지 확인
            for text in texts:
                if text in shape.text:
                    shapes_with_texts.append(shape)
                    break
    return shapes_with_texts


def insert_signalment_to_slide(prs, slide_index, signalment_text):
    # 슬라이드를 가져옵니다.
    slide = prs.slides[slide_index]

    # 텍스트가 삽입될 도형을 찾습니다.
    shape = None
    for shp in slide.shapes:
        if shp.has_text_frame:
            shape = shp
            break
    
    if shape:
        # 기존 텍스트를 지우고 새로운 텍스트를 삽입합니다.
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = signalment_text
        p.font.size = Pt(24)  # 글꼴 크기를 설정합니다.



#슬라이드 복사
def copy_slide(prs, slide_index):
    slide = prs.slides[slide_index]  # 인덱스에 해당하는 슬라이드를 가져옴
    slide_layout = slide.slide_layout  # 슬라이드 레이아웃을 가져옴
    new_slide = prs.slides.add_slide(slide_layout)  # 동일한 레이아웃으로 새로운 슬라이드를 추가

    for shape in slide.shapes:  # 원본 슬라이드의 모든 도형을 순회
        if isinstance(shape, (GroupShape, Connector)):
            continue  # 그룹 도형이나 커넥터는 복사하지 않음

        if not shape.has_text_frame:  # 텍스트 프레임이 없는 도형인 경우
            if shape.shape_type == 13:  # 도형이 이미지인 경우
                image_path = f'extracted_images/temp_image_{shape.shape_id}.png'
                with open(image_path, 'wb') as f:
                    f.write(shape.image.blob)  # 이미지 데이터를 파일로 저장
                     # 이미지를 새로운 슬라이드에 추가
                new_slide.shapes.add_picture(image_path, shape.left, shape.top, shape.width, shape.height) 
            else:
                # 도형을 새로운 슬라이드에 추가
                new_shape = new_slide.shapes.add_shape(shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height)  
        else:
             # 텍스트 상자를 새로운 슬라이드에 추가
            new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
            # 원본 도형의 모든 문단을 순회 
            for paragraph in shape.text_frame.paragraphs:  
                # 새로운 문단을 추가
                new_paragraph = new_shape.text_frame.add_paragraph() 
                # 원본 문단의 텍스트를 복사 
                new_paragraph.text = paragraph.text  

    return new_slide  # 복사된 슬라이드를 반환


#vital check 표 추가
def add_vital_check_table_to_slide(prs, headers, data_list, slide_layout_index):
    if data_list is None or (isinstance(data_list, pd.DataFrame) and data_list.empty) or (isinstance(data_list, list) and not data_list):
        #데이터가 없으면 빈 딕셔너리 추가
        data_list = [{}]
    #지정된 인덱스의 슬라이드 레이아웃 가져옴
    slide_layout = prs.slide_layouts[slide_layout_index]
    #새로운 슬라이드 추가
    slide = prs.slides.add_slide(slide_layout)
    #슬라이드 제목을 가져옴
    title = slide.shapes.title
    if title:
        title.text = "Vital Check"
    else:
        #텍스트 상자 추가
        txBox = slide.shapes.add_textbox(Cm(2), Cm(1), Cm(6), Cm(1.5))
        tf = txBox.text_frame
        tf.text = "Vital Check"

    rows = len(data_list) + 1
    cols = len(headers)
    
    # 슬라이드 크기 가져오기
# 슬라이드 크기 가져오기
    slide_width = prs.slide_width
    slide_height = prs.slide_height
        
        # 표의 너비와 높이 설정 (필요에 따라 조정 가능)
    table_width = Inches(6)  # 예시로 설정한 표 너비
    table_height = Inches(3)  # 예시로 설정한 표 높이
        
        # 표의 왼쪽 및 상단 위치를 슬라이드 중앙으로 설정
    left = (slide_width - table_width) / 2
    top = (slide_height - table_height) / 2
        
        # 표를 슬라이드 중앙에 추가
    table = slide.shapes.add_table(rows, cols, left, top, table_width, table_height).table

    # 헤더를 표에 추가
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                # 헤더 텍스트의 폰트 크기를 5pt로 설정
                run.font.size = Pt(5)
# 데이터를 표에 추가
    for row_idx, data in enumerate(data_list):
        for col_idx, header in enumerate(headers):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(data.get(header, ''))
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(5)

def add_lab_table_by_date_to_slide(prs, lab_df, slide_layout_index):
    # 날짜별로 그룹화하여 슬라이드에 표 추가
    grouped = lab_df.groupby('date')
    
    for date, group in grouped:
        slide_layout = prs.slide_layouts[slide_layout_index]
        slide = prs.slides.add_slide(slide_layout)
        
        # 슬라이드 제목 설정
        title = slide.shapes.title
        if title:
            title.text = f"Lab Results - {date}"
        else:
            txBox = slide.shapes.add_textbox(Cm(2), Cm(1), Cm(6), Cm(1.5))
            tf = txBox.text_frame
            tf.text = f"Lab Results - {date}"
        
        # 데이터와 헤더 설정
        rows = len(group) + 1
        cols = len(group.columns) - 1  # date 컬럼은 슬라이드 제목으로 대체하므로 제외
        
        # 슬라이드 크기 가져오기
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # 표의 너비와 높이 설정 (필요에 따라 조정 가능)
        table_width = Inches(6)  # 예시로 설정한 표 너비
        table_height = Inches(3)  # 예시로 설정한 표 높이
        
        # 표의 왼쪽 및 상단 위치를 슬라이드 중앙으로 설정
        left = (slide_width - table_width) / 2
        top = (slide_height - table_height) / 2
        
        # 표를 슬라이드 중앙에 추가
        table = slide.shapes.add_table(rows, cols, left, top, table_width, table_height).table
        
        # 헤더 설정
        headers = ['검사명', '결과값', '단위', 'MIN', 'MAX', 'Description', 'status']
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(5)
        
        # 데이터 추가
        for row_idx, (_, data) in enumerate(group.iterrows()):
            for col_idx, header in enumerate(headers):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(data[header])
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(5)



def add_new_slide_with_chart_image(prs, image1, image2=None):
    slide_layout = None
    # 프레젠테이션의 슬라이드 레이아웃을 순회
    for layout in prs.slide_layouts:  
        # 'Title Only' 레이아웃을 찾음
        if layout.name == 'Title Only':  
            slide_layout = layout
            break
# 'Title Only' 레이아웃을 찾지 못한 경우
    if not slide_layout:  
        print("Title Only layout not found. Using default layout.")
        # 기본 레이아웃을 사용
        slide_layout = prs.slide_layouts[0]  
# 새로운 슬라이드를 추가
    new_slide = prs.slides.add_slide(slide_layout) 
# 슬라이드의 제목을 가져옴
    title = new_slide.shapes.title  
    # 제목이 있으면 "Chart Image"로 설정
    if title:
        title.text = "Chart Image"  
    else:
        txBox = new_slide.shapes.add_textbox(Cm(2), Cm(1), Cm(6), Cm(1.5)) 
        tf = txBox.text_frame
        # 텍스트 상자에 "Chart Image" 텍스트를 추가
        tf.text = "Chart Image"  
# 임시 파일 생성
    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_image1:  
        # 첫 번째 이미지를 저장
        image1.save(tmp_image1, format="PNG")  
        tmp_image1_path = tmp_image1.name
        print(f"Saving image1 to {tmp_image1_path}")
# 임시 파일이 존재하는지 확인
    if os.path.exists(tmp_image1_path):  
        try:
            # 이미지를 열기
            image1 = Image.open(tmp_image1_path)  
            width, height = image1.size
            max_width = Inches(5)
            max_height = Inches(4)
            # 크기 비율 계산
            ratio = min(max_width / width, max_height / height)  
            new_width = width * ratio
            new_height = height * ratio
        # 슬라이드에 이미지 추가
            new_slide.shapes.add_picture(tmp_image1_path, Inches(1), Inches(2), width=new_width, height=new_height)  
            print("Image1 added to slide successfully.")
        except Exception as e:
            print(f"Error adding image1 to slide: {e}")

    if image2:  # 두 번째 이미지가 있는 경우
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_image2:  # 임시 파일 생성
            image2.save(tmp_image2, format="PNG")  # 두 번째 이미지를 저장
            tmp_image2_path = tmp_image2.name
            print(f"Saving image2 to {tmp_image2_path}")
        # 임시 파일이 존재하는지 확인
        if os.path.exists(tmp_image2_path): 
            # 이미지를 열기
            try:
                image2 = Image.open(tmp_image2_path)  
                width, height = image2.size
                max_width = Inches(5)
                max_height = Inches(4)
                 # 크기 비율 계산
                ratio = min(max_width / width, max_height / height) 
                new_width = width * ratio
                new_height = height * ratio
            # 슬라이드에 이미지 추가
                new_slide.shapes.add_picture(tmp_image2_path, Inches(6), Inches(2), width=new_width, height=new_height)  
                print("Image2 added to slide successfully.")
            except Exception as e:
                print(f"Error adding image2 to slide: {e}")

def plot_vital_check_graph(vital_check_data_list):
    # vital_check_data_list에서 날짜와 BW(Kg) 추출
    dates = vital_check_data_list["날짜"].tolist()
    bw_values = vital_check_data_list["BW(Kg)"].tolist()

    # 추출된 데이터 출력
    print("Dates:", dates)
    print("BW Values:", bw_values)
    
    if not dates or not bw_values:
        print("날짜와 BW(Kg) 값을 찾을 수 없습니다.")
        return

    #bw 데이터 표시
    bw_values = [float(value) for value in bw_values]

    # x값 정렬
    indices = list(range(1, len(dates) + 1))

    # 그래
    plt.figure(figsize=(12, 6))
    plt.plot(indices, bw_values, marker='o', label='BW (Kg)')
    plt.xlabel('Measurement Index')
    plt.ylabel('BW(Kg)')
    plt.ylim(0, 4)
    plt.title('BW(Kg) Over Time')
    plt.legend()
    plt.grid(True)

    # Set the x-axis ticks to be the indices
    plt.xticks(indices)

    # 이미지 저장
    graph_path = 'vital_check_graph.png'
    plt.tight_layout()
    plt.savefig(graph_path)
    plt.close()

    return graph_path
# pdf 추출
def get_ppt_input(pdf_path):
    pages = load_pdf(pdf_path)

    soap = get_soap(pages)
    info = get_animal_data(pages)

    llm_input = get_LLM_input(info, soap)
    #print(llm_input)

    #------------------------get LLM input--------------------------#
    from langchain_openai import ChatOpenAI
    llm = ChatOpenAI(model_name='gpt-4o', temperature=0)



    from prompt.prompt import get_prompt, get_parser
    print("1")
    prompt = get_prompt()
    print("2")
    parser = get_parser()

    #------------------------chain--------------------------#
    from langchain_core.runnables import RunnablePassthrough
    print("3")

    print(prompt)
    chain = {'input' : RunnablePassthrough()} | prompt | llm | parser

    print("4")
    print(llm_input)
    answer = chain.invoke(llm_input)
    #print(answer)
    return answer


def extract_signalment_data(signalment_str):
    # 텍스트를 줄 단위로 분리한 다음, 쉼표 단위로 분리
    parts = [item.strip() for line in signalment_str.split('\n') for item in line.split(',')]

    # 데이터 추출을 위한 키워드 목록
    keywords = ['Name', 'Breed', 'Species', 'Age', 'Sex', 'BW']
    data = {}

    # 키워드를 기준으로 텍스트 분리
    for part in parts:
        for keyword in keywords:
            if part.startswith(keyword):
                data[keyword] = part.split(': ')[1].strip()
                break

    return data

def insert_signalment_to_slide_title_and_body(prs, slide_index, signalment_text, append=True):
    """
    슬라이드의 제목과 본문(텍스트 상자)에 텍스트를 추가하는 함수.
    제목에는 'signalment' 텍스트를 추가하고, 본문에는 signalment_text를 추가합니다.
    append가 True인 경우 본문 텍스트를 기존 텍스트 뒤에 추가하고, False인 경우 기존 텍스트를 덮어씁니다.
    """
    # 슬라이드를 가져옵니다.
    slide = prs.slides[slide_index]

    # 제목 텍스트 상자를 찾고 "signalment"를 삽입합니다.
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "signalment"

    # 본문(텍스트 상자)을 찾습니다.
    body_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape != title_shape:
            body_shape = shape
            break
    
    if body_shape:
        text_frame = body_shape.text_frame
        if not append:
            text_frame.clear()  # 기존 텍스트 지우기
        p = text_frame.add_paragraph()
        p.text = signalment_text
        p.font.size = Pt(13)            # 글꼴 크기 설정
        p.font.bold = False             # 글꼴 굵기 설정
        p.font.name = 'Arial'           # 글꼴 설정
        
        body_shape.left = Inches(1)     # 왼쪽 여백 설정
        body_shape.top = Inches(1.5)      # 상단 여백 설정
        body_shape.width = Inches(8)    # 너비 설정
        body_shape.height = Inches(4) 
    else:
        print(f"Slide {slide_index} does not have a content area.")

def insert_text_to_slide(prs, slide_index, title_text, body_text):
    """
    슬라이드에 제목과 본문 텍스트를 삽입하는 함수.
    """
    # 슬라이드를 가져옵니다.
    slide = prs.slides[slide_index]

    # 제목 텍스트 상자에 텍스트 삽입
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = title_text

    # 본문(텍스트 상자)을 찾습니다.
    body_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape != title_shape:
            body_shape = shape
            break

    if body_shape:
        text_frame = body_shape.text_frame
        text_frame.clear()  # 기존 텍스트 지우기
        p = text_frame.add_paragraph()
        run = p.add_run()  # 개별 실행(run) 추가
        run.text = body_text

        # 폰트 스타일 설정
        run.font.size = Pt(13)            # 글꼴 크기 설정
        run.font.bold = False             # 글꼴 굵기 설정
        run.font.name = 'Noto Sans KR'    # 글꼴 설정
        
         # 텍스트 상자 기본 서식 무시하기 위해 스타일 초기화 시도
        p.font.size = Pt(13)  # 단락 수준에서 폰트 크기를 재설정

        # 마지막 슬라이드(전체 8번째 슬라이드)를 위한 특별한 여백 및 크기 조정
       
            # 기본 위치와 크기 조정
        body_shape.left = Inches(1)     # 왼쪽 여백 설정
        body_shape.top = Inches(1.5)      # 상단 여백 설정
        body_shape.width = Inches(8)    # 너비 설정
        body_shape.height = Inches(4)  # 높이 설정
    else:
        print(f"Slide {slide_index} does not have a content area.")
        
        
def get_pptx_tables(pdf_path):
    pages = load_pdf(pdf_path)

    lab_tables = get_lab(pages)
    # HCT, WBC, PLT
    extracted_lab_tables = []

    for data in lab_tables:
        date = data['date']
        table = data['table']
        

        hct_row = table[table['검사명'] == 'HCT']
        wbc_row = table[table['검사명'] == 'WBC']
        plt_row = table[table['검사명'] == 'PLT']

        if not hct_row.empty:
            hct_row = hct_row.copy()
            hct_row['date'] = date
            extracted_lab_tables.append(hct_row)
            # print("date : ", date)
            # print("hct : ",hct_row)
        if not wbc_row.empty:
            wbc_row = wbc_row.copy()
            wbc_row['date'] = date
            extracted_lab_tables.append(wbc_row)
            # print("date : ", date)
            # print(wbc_row)
        if not plt_row.empty:
            plt_row = plt_row.copy()
            plt_row['date'] = date
            extracted_lab_tables.append(plt_row)
            # print("date : ", date)
            # print(plt_row)

    result_df = pd.concat(extracted_lab_tables, ignore_index=True)
    print(result_df)
    return result_df

def add_textbox_to_slides(prs, text=" fff"):
    """
    첫 번째와 마지막 슬라이드를 제외한 모든 슬라이드에 텍스트 박스를 추가하는 함수.
    """
    slide_count = len(prs.slides)  # 전체 슬라이드 개수

    # 첫 번째와 마지막 슬라이드를 제외하고 순회
    for slide_index in range(1, slide_count - 1):
        slide = prs.slides[slide_index]

        # 텍스트 상자 추가 (위치와 크기는 필요에 따라 조정)
        left = Inches(1)
        top = Inches(5)
        width = Inches(8)
        height = Inches(1)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.clear()

        # 텍스트 추가
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(13)  # 폰트 크기 설정
        p.font.bold = True  # 굵게 설정

def main_presentation():
    pdf_path = "/Users/sinequanon/Documents/PerfectS/data/pdf_data/TEST_01.pdf"
    soap_result, vital_check_data_list, lab_data_list, animaldata = main(pdf_path)
    pptx_file_path = '/Users/sinequanon/Documents/PerfectS/data/pptx/template3.pptx'
    
    
    if not os.path.exists(pptx_file_path):
        print("PPTX file does not exist.")
        return

    prs = Presentation(pptx_file_path)
    answer = get_ppt_input(pdf_path)
    print("받은 데이터===========")
    print(answer)
    result_df = get_pptx_tables(pdf_path)

    
    # signalment 부분 추출 및 처리
    signalment_str = answer.signalment

    # 문자열에서 데이터를 추출하여 딕셔너리로 반환
    extracted_data = extract_signalment_data(signalment_str)

    # Signalment 텍스트 생성
    signalment_text = (
        f"이름: {extracted_data.get('Name', 'N/A')}\n"
        f"품종: {extracted_data.get('Breed', 'N/A')}\n"
        f"성별: {extracted_data.get('Sex', 'N/A')}\n"
        f"종류: {extracted_data.get('Species', 'N/A')}\n"
        f"나이: {extracted_data.get('Age', 'N/A')}\n"
        f"현재 무게: {extracted_data.get('BW', 'N/A')}\n"
    )

    # 두 번째 슬라이드에 signalment 삽입
    insert_signalment_to_slide_title_and_body(prs, slide_index=1, signalment_text=signalment_text, append=True)
 # answer 문자열을 파싱하여 key-value 쌍을 얻음
 # MedicalSummary 객체의 속성 추출
    attributes = ['chief_complaint', 'body_check', 'diagnosis', 'plan_edu', 'Sugery', 'a_record', 'postcare']
    
    # 첫 번째 항목은 두 번째 슬라이드에 삽입
    slide_index = 2

    for attr in attributes:
        if hasattr(answer, attr):
            title = attr
            body = getattr(answer, attr)

            # 슬라이드에 텍스트 삽입
            insert_text_to_slide(prs, slide_index, title, body)

            # 두 번째 항목부터는 슬라이드를 복제하여 사용
            if slide_index == 1:
                slide_index += 1
            else:
                slide_layout = prs.slides[1].slide_layout
                new_slide = prs.slides.add_slide(slide_layout)
                slide_index = len(prs.slides) - 1
    # 데이터를 불러오기
    

    xml_slides = prs.slides._sldIdLst  
    slides = list(xml_slides)
    for i in [3, 2]:  # 뒤에서부터 제거
        xml_slides.remove(slides[i])
        
    """======================================================= SOAP PAGES DONE ======================================================="""





























    # 이미지 추출
    images = extract_images_from_pdf(pdf_path)
    get_ppt_input(pdf_path)

    for i in range(0, len(images), 2):
        add_new_slide_with_chart_image(prs, images[i], images[i + 1] if i + 1 < len(images) else None)
    
    # Ensure any remaining images are added
    if len(images) % 2 != 0:
        add_new_slide_with_chart_image(prs, images[-1])


    """======================================================= CHART PAGES DONE ======================================================="""










    # Vital Check 그래프 생성
    graph_path = plot_vital_check_graph(vital_check_data_list)

    # 그래프 슬라이드 추가
    if graph_path:
        # "Title Only" 레이아웃 사용하여 슬라이드 추가
        slide_layout = None
        for layout in prs.slide_layouts:
            if layout.name == 'Title Only':
                slide_layout = layout
                break

        if not slide_layout:
            slide_layout = prs.slide_layouts[0]

        graph_slide = prs.slides.add_slide(slide_layout)
        title = graph_slide.shapes.title
        if title:
            title.text = "Vital Check Graph"

        graph_slide.shapes.add_picture(graph_path, Inches(1), Inches(2), width=Inches(8), height=Inches(4.5))

    slide_layout_index = None
    for i, layout in enumerate(prs.slide_layouts):
        if layout.name == 'Title Only':
            slide_layout_index = i
            break

    if slide_layout_index is None:
        slide_layout_index = 0  # Default to the first slide layout if none is found

    vital_headers = ['날짜', '시간', 'BW(Kg)', 'BT(C)', 'BP(mmHg)', 'HR(/min)', 'Sign']
    lab_headers = ['검사명', '결과값', '단위', 'MIN', 'MAX', 'Description']

    add_vital_check_table_to_slide(prs, vital_headers, vital_check_data_list.to_dict('records'), slide_layout_index)
    add_lab_table_by_date_to_slide(prs, result_df, slide_layout_index)


    slide_layout = None
    for layout in prs.slide_layouts:
            if layout.name == 'Centered Text':
                slide_layout = layout
                break

    if not slide_layout:
        print("Centered Text layout not found. Using default layout.")
        slide_layout = prs.slide_layouts[0]

    thank_you_slide = prs.slides.add_slide(slide_layout)

    # 기존 텍스트 상자를 찾기
    textbox = None
    for shape in thank_you_slide.shapes:
        if shape.has_text_frame:  # 텍스트 상자가 있는지 확인
            textbox = shape
            break

    if textbox:
        # 텍스트 상자가 있는 경우 기존 내용을 지우지 않고 텍스트 추가
        text_frame = textbox.text_frame

        # 단락 추가
        p = text_frame.add_paragraph()
        p.text = "감사합니다"
        p.font.size = Pt(40)
        p.font.bold = True

        # 텍스트의 가로 정렬을 중앙으로 설정
        p.alignment = PP_ALIGN.CENTER

        # 텍스트 상자의 세로 정렬을 중앙으로 설정
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    else:
        print("본문 텍스트 상자를 찾을 수 없습니다.")


    xml_slides = prs.slides._sldIdLst  
    slides = list(xml_slides)
    for i in [9, 8]:  # 뒤에서부터 제거 (슬라이드 인덱스는 0부터 시작하므로 8, 9, 10은 7, 8, 9에 해당)
        xml_slides.remove(slides[i])
        
    add_textbox_to_slides(prs, text=" ")
    
    output_path = '/Users/sinequanon/Documents/PerfectS/data/pptx/test_pptx.pptx'
    prs.save(output_path)

if __name__ == "__main__":
    main_presentation()
