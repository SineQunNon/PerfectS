from pptx.util import Inches, Cm, Pt
from pptx import Presentation
import logging

from Rag.llm.get_ppt_input import get_ppt_input

def get_signalment_page_data(signalment_str):
    try:
        parts = [item.strip() for line in signalment_str.split('\n') for item in line.split(',')]

        keywords = ['Name', 'Breed', 'Species', 'Age', 'Sex', 'BW']
        data = {}

        for part in parts:
            for keyword in keywords:
                if part.startswith(keyword):
                    data[keyword] = part.split(': ')[1].strip()

        # 텍스트 생성
        signalment_text = (
            f"• Name : {data.get('Name', 'N/A')}\n"
            f"• Breed : {data.get('Breed', 'N/A')}\n"
            f"• Species : {data.get('Species', 'N/A')}\n"
            f"• Age : {data.get('Age', 'N/A')}\n"
            f"• Sex : {data.get('Sex', 'N/A')}\n"
            f"• BW : {data.get('BW', 'N/A')}\n"
        )
    
        return signalment_text
    except Exception as e:
        logging.error(f"Error processing signalment in get_signalment_page_data : {e}")
        return "N/A"

def get_soap_text(body_text):
    body_texts = [item.strip() for line in body_text.split('\n') for item in line.split(',')]
    text = ""

    for txt in body_texts:
        text += f"• {txt}\n"

    return text



def add_soap_pages(prs, slide_index, title_text, body_text, num, append=True):
    """
    슬라이드의 제목과 본문(텍스트 상자)에 텍스트를 추가하는 함수.
    제목에는 'signalment' 텍스트를 추가하고, 본문에는 signalment_text를 추가합니다.
    append가 True인 경우 본문 텍스트를 기존 텍스트 뒤에 추가하고, False인 경우 기존 텍스트를 덮어씁니다.
    """
    # 슬라이드를 가져옵니다.
    prs.slides.add_slide(prs.slide_layouts[1])
    slide = prs.slides[slide_index]

    titlenum_shape = slide.shapes[0]
    if num < 10:
        titlenum_shape.text = "0" + str(num)
    else:
        titlenum_shape.text = str(num)
    # 제목 텍스트 상자를 찾고 "signalment"를 삽입합니다.
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = title_text

    # 본문(텍스트 상자)을 찾습니다.
    body_shape = slide.shapes[-1]
    body_shape.text = body_text

def chunk_list(lst, n):
    for i in range(0, len(lst), n):
        yield lst[i:i + n]

def make_soap_pages(prs, pdf_path):
    answer = get_ppt_input(pdf_path)
    signalment_text = get_signalment_page_data(answer.signalment)

    add_soap_pages(prs, slide_index=1, title_text="Signalment", body_text=signalment_text, num=1,append=True, )


    attributes = ['chief_complaint', 'body_check', 'diagnosis', 'plan_edu', 'surgery', 'a_record', 'postcare']
    slide_index = 2
    t_num = 2
    for attr in attributes:
        if hasattr(answer, attr):
            if attr == 'chief_complaint':
                title = 'Chief Complaint'
            elif attr == 'body_check':
                title = 'Body Check'
            elif attr == 'diagnosis':
                title = 'Diagnosis'
            elif attr == 'plan_edu':
                title = 'Plan Education'
            elif attr == 'surgery':
                title = 'Surgery'
            elif attr == 'a_record':
                title = 'Anesthesia Record'
            elif attr == 'postcare':
                title = 'Postoperative Care'
            else:
                title = 'N/A'
            
            body = getattr(answer, attr)
            
            body = get_soap_text(body)

            if len(body.split('\n')) < 7:
                add_soap_pages(prs, slide_index, title, body, t_num, append=True)
                slide_index += 1
            else:
                body_lines = body.split('\n')

                for chunk in chunk_list(body_lines, 6):
                    chunk_text = "\n".join(chunk)
                    add_soap_pages(prs, slide_index, title, chunk_text, t_num, append=True)
                    slide_index+=1
            t_num += 1