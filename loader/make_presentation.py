import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from main import *

#슬라이드에서 특정 도형rhk 택스트를 선택한다
def select_shape_by_text(slide, text):
    for shape in slide.shapes:
        if shape.has_text_frame and text in shape.text:
            return shape
    return None

#여러개의 텍스트 일치 비교
def select_shapes_by_texts(slide, texts):
    shapes_with_texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for text in texts:
                if text in shape.text:
                    shapes_with_texts.append(shape)
                    break  # 한 텍스트가 일치하면 다른 텍스트는 확인할 필요 없음
    return shapes_with_texts


def insert_data(slide, animal_data):
    signalment_shape = select_shape_by_text(slide, "Signalment_data")
    if signalment_shape:
        signalment_shape.text_frame.clear()
        signalment_text = (
            f"이름: {animal_data['name']}\n"
            f"품종: {animal_data['breed']}\n"
            f"성별: {animal_data['sex']}\n"
            f"종류: {animal_data['species']}\n"
            f"나이: {animal_data['birth']}\n"
            f"현재 무게: {animal_data['weight']}\n"
            f"보호자: {animal_data['protector']}"
        )
        signalment_shape.text_frame.text = signalment_text
def insert_soap(slide, soap_data):
    # Subjective 데이터 삽입
    subjective_shape = select_shape_by_text(slide, 'Subjective text')
    if subjective_shape:
        
        subjective_shape.text_frame.clear()
        #Subjective text에 추출된 데이터를 넣는다
        for paragraph in soap_data.get('subjective', []):
            p = subjective_shape.text_frame.add_paragraph()
            p.text = paragraph

    else:
        print("Subjective shape not found")

    # Objective 데이터 삽입
    objective_shape = select_shape_by_text(slide, 'Objective text')
    if objective_shape:
    #Objective text에 추출된 데이터를 넣는다
        objective_shape.text_frame.clear()
        for paragraph in soap_data.get('objective', []):
            p = subjective_shape.text_frame.add_paragraph()
            p.text = paragraph
    else:
        print("Objective shape not found")

    # Assessment 데이터 삽입
    assessment_shape = select_shape_by_text(slide, 'Assessment text')
    if assessment_shape:
        #Assessment text에 추출된 데이터를 넣는다
        assessment_shape.text_frame.clear()
        for paragraph in soap_data.get('assessment', []):
            p = assessment_shape.text_frame.add_paragraph()
            p.text = paragraph
    
    else:
        print("Assessment shape not found")

    # Plan 데이터 삽입
    plan_shape = select_shape_by_text(slide, 'Plan text')
    if plan_shape:
    #Plan text에 추출된 데이터를 넣는다
        plan_shape.text_frame.clear()
        for paragraph in soap_data.get('plan', []):
            p = plan_shape.text_frame.add_paragraph()
            p.text = paragraph

    else:
        print("Plan shape not found")
#특정 슬라이드 복사
def copy_slide(prs, slide_index):
    #prs - 프레젠테이션 , 
    #복사할 슬라이드와 레이아웃을 가져온다
    slide = prs.slides[slide_index]
    slide_layout = slide.slide_layout
    #새 슬라이드를 만든다
    new_slide = prs.slides.add_slide(slide_layout)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            if shape.shape_type == 13: #도형이 이미지인지 확인
                #이미지 도형인 경우 이미지 파일로 저장한다
                image_path = '/Users/sinequanon/Documents/PerfectS/data/images/temp.png'
                with open(image_path, 'wb') as f:
                    f.write(shape.image.blob)
                    #텍스프 프레임이 있을 경우 동일한 위치와 크기로 새 슬라이드에 추가
                new_slide.shapes.add_picture(image_path, shape.left, shape.top, shape.width, shape.height)
            else:
                #이미지가 아닌 도형 복사
                new_shape = new_slide.shapes.add_shape(shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height)
        else:
            #텍스트 프레임이 있는 도형 복사
            new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
            for paragraph in shape.text_frame.paragraphs:
                new_paragraph = new_shape.text_frame.add_paragraph()
                new_paragraph.text = paragraph.text

    return new_slide# 새로운 슬라이드 반환

#vital check 표 넣기
def add_vital_check_table_to_slide(prs, headers, data_list):
    if data_list is None or (isinstance(data_list, pd.DataFrame) and data_list.empty) or (isinstance(data_list, list) and not data_list):
        data_list = [{}]

    slide_layout = prs.slide_layouts[5]  # 빈 슬라이드 레이아웃 추가
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Vital Check"

    rows = len(data_list) + 1  # 헤더를 포함한 행 수
    cols = len(headers)   # 헤더의 수
    #테이블 추가
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(5.0)).table
    #헤더 채우기
    for col_idx, header in enumerate(headers):
        table.cell(0, col_idx).text = header

    #데이터 채우기
    for row_idx, data in enumerate(data_list):
        for col_idx, header in enumerate(headers):
            table.cell(row_idx + 1, col_idx).text = str(data[header]) if isinstance(data, dict) else str(data)

#lab Result 채우기
def add_lab_table_to_slide(prs, headers, lab_data_list):
    if lab_data_list is None or (isinstance(lab_data_list, pd.DataFrame) and lab_data_list.empty) or (isinstance(lab_data_list, list) and not lab_data_list):
        lab_data_list = [{}]
    
    for lab_data in lab_data_list:
        if isinstance(lab_data, dict) and isinstance(lab_data.get('table'), pd.DataFrame):
            slide_layout = prs.slide_layouts[5]  #빈 슬라이드에 레이아웃 추가
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = f"Lab Results - {lab_data.get('date', '')}"
            
            rows = len(lab_data['table']) + 1 #헤더를 포함한 행수
            cols = len(headers) #헤더의 수
            table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(5), Inches(3.0)).table
            #헤더 채우기
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = header
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(5) 
            #데이터 채우기
            for row_idx, (_, data) in enumerate(lab_data['table'].iterrows()):
                for col_idx, header in enumerate(headers):
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = str(data[header])
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(5) 

def main_presentation():
    # PDF 파일 경로 설정
    pdf_path = "/Users/sinequanon/Documents/PerfectS/data/pdf_data/TEST_05.pdf"
    # 변수 들고오기
    soap_result, vital_check_data_list, lab_data_list, animal_data = main(pdf_path)
    print(soap_result, vital_check_data_list, lab_data_list, animal_data)

    
    pptx_file_path = '/Users/sinequanon/Documents/PerfectS/data/template.pptx'
    if not os.path.exists(pptx_file_path):
        print("PPTX file does not exist.")
        return

    prs = Presentation(pptx_file_path)

    # 이미지가 있는 경로
    image_directory = '/Users/sinequanon/Documents/PerfectS/data/images'
    images = [os.path.join(image_directory, f) for f in os.listdir(image_directory) if os.path.isfile(os.path.join(image_directory, f)) and f.endswith(('png', 'jpg', 'jpeg', 'gif', 'bmp'))]


    #특정 텍스트를 포함하는 슬라이드 찾기
    slide_layout = None
    slide_index = None
    for i, slide in enumerate(prs.slides):
        slide_title_shape = select_shape_by_text(slide, "Chart Image")
        if slide_title_shape:
            slide_layout = slide.slide_layout
            slide_index = i
            break

#이미지 추가
    if slide_layout:
        slide = prs.slides[slide_index]
        for i in range(0, len(images), 2):
            if i != 0:
                slide = copy_slide(prs, slide_index)
            slide.shapes.add_picture(images[i], Cm(2), Cm(5), width=Cm(10), height=Cm(7))
            if i + 1 < len(images):
                slide.shapes.add_picture(images[i + 1], Cm(14), Cm(5), width=Cm(10), height=Cm(7))
#테이블 추가하기
    vital_headers = ['날짜', '시간', 'BW(Kg)', 'BT(C)', 'BP(mmHg)', 'HR(/min)', 'Sign']
    lab_headers = ['검사명', '결과값', '단위', 'MIN', 'MAX', 'Description']

    add_vital_check_table_to_slide(prs, vital_headers, vital_check_data_list.to_dict('records'))
    add_lab_table_to_slide(prs, lab_headers, lab_data_list)

#signalment_data가 포함되어 있으면 animal 데이터를 넣는다
    for slide in prs.slides:
        animal_slide = select_shape_by_text(slide, "Signalment_data")
        if animal_slide:
            insert_data(slide, animal_data)


    #해당 텍스트가 포함되어 있으면 soap_data를 넣는다         
    texts_to_find = ["Subjective text", "Objective text", "Assessment text", "Plan text"]
    for slide in prs.slides:
        soap_shapes = select_shapes_by_texts(slide, texts_to_find)
        for soap_shap in soap_shapes:
            for soap_data in soap_result:
                insert_data(slide, soap_data)
                
    output_path = '/Users/sinequanon/Documents/PerfectS/data/pptx/test05.pptx'
    prs.save(output_path)

if __name__ == "__main__":
    main_presentation()

